"""Document loading and text extraction."""

import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional, Dict, Any

from signalinator_core import get_logger

logger = get_logger(__name__)


@dataclass
class LoadedDocument:
    """A loaded document with extracted text."""

    content: str
    filename: str
    file_path: str
    document_type: str
    page_count: int = 1
    metadata: Dict[str, Any] = field(default_factory=dict)


class DocumentLoader:
    """Multi-format document loader."""

    SUPPORTED_EXTENSIONS = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "docx",
        ".pptx": "pptx",
        ".ppt": "pptx",
        ".odt": "odt",
        ".txt": "txt",
        ".md": "txt",
        ".csv": "txt",
        ".json": "txt",
        ".yaml": "txt",
        ".yml": "txt",
        ".log": "txt",
    }

    def load(self, file_path: str) -> Optional[LoadedDocument]:
        """Load a document from file path."""
        path = Path(file_path)
        if not path.exists():
            logger.error(f"File not found: {file_path}")
            return None

        ext = path.suffix.lower()
        doc_type = self.SUPPORTED_EXTENSIONS.get(ext)

        if not doc_type:
            logger.warning(f"Unsupported file type: {ext}")
            return None

        try:
            if doc_type == "pdf":
                return self._load_pdf(path)
            elif doc_type == "docx":
                return self._load_docx(path)
            elif doc_type == "pptx":
                return self._load_pptx(path)
            elif doc_type == "odt":
                return self._load_odt(path)
            else:
                return self._load_text(path)

        except Exception as e:
            logger.error(f"Error loading {file_path}: {e}")
            return None

    def _load_pdf(self, path: Path) -> LoadedDocument:
        """Load PDF document."""
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = []

        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[Page {i}]\n{text}")

        return LoadedDocument(
            content="\n\n".join(pages),
            filename=path.name,
            file_path=str(path),
            document_type="pdf",
            page_count=len(reader.pages),
        )

    def _load_docx(self, path: Path) -> LoadedDocument:
        """Load Word document."""
        from docx import Document

        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        return LoadedDocument(
            content="\n\n".join(paragraphs),
            filename=path.name,
            file_path=str(path),
            document_type="docx",
            page_count=1,
        )

    def _load_pptx(self, path: Path) -> LoadedDocument:
        """Load PowerPoint document."""
        from pptx import Presentation

        prs = Presentation(str(path))
        slides = []

        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))

        return LoadedDocument(
            content="\n\n".join(slides),
            filename=path.name,
            file_path=str(path),
            document_type="pptx",
            page_count=len(prs.slides),
        )

    def _load_odt(self, path: Path) -> LoadedDocument:
        """Load OpenDocument text."""
        from odf import text as odf_text
        from odf.opendocument import load as odf_load

        doc = odf_load(str(path))
        paragraphs = []

        for p in doc.getElementsByType(odf_text.P):
            text = "".join(
                node.data for node in p.childNodes
                if hasattr(node, "data")
            )
            if text.strip():
                paragraphs.append(text)

        return LoadedDocument(
            content="\n\n".join(paragraphs),
            filename=path.name,
            file_path=str(path),
            document_type="odt",
            page_count=1,
        )

    def _load_text(self, path: Path) -> LoadedDocument:
        """Load plain text document."""
        content = path.read_text(encoding="utf-8", errors="ignore")

        return LoadedDocument(
            content=content,
            filename=path.name,
            file_path=str(path),
            document_type="txt",
            page_count=1,
        )

    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """Check if file type is supported."""
        ext = Path(file_path).suffix.lower()
        return ext in cls.SUPPORTED_EXTENSIONS
